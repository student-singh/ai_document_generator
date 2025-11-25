# pptx_generator.py
from pptx import Presentation

def create_pptx(project):
    prs = Presentation()
    for i, title in enumerate(project["structure"]):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = project["content"].get(str(i), {}).get("text", "")
    path = f"uploads/{project['_id']}.pptx"
    prs.save(path)
    return path